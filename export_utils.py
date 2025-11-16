# export_utils.py
# -*- coding: utf-8 -*-
"""
通用导出工具：
- TXT
- Word（.docx）
- PPT（.pptx）
- Excel（.xlsx）

设计思路：
1. 不强依赖第三方库：如果缺少 python-docx / python-pptx / openpyxl，会给出友好提示。
2. 统一用 UTF-8，文件放到 ./exports 目录下。
"""

import os
from datetime import datetime
from typing import List, Dict


# ---------- 公共工具 ----------

def ensure_export_dir() -> str:
    export_dir = os.path.join(os.getcwd(), "exports")
    if not os.path.exists(export_dir):
        os.makedirs(export_dir, exist_ok=True)
    return export_dir


def timestamp_str() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


# ---------- TXT 导出 ----------

def export_to_txt(basename: str, content: str) -> str:
    export_dir = ensure_export_dir()
    filename = f"{basename}_{timestamp_str()}.txt"
    path = os.path.join(export_dir, filename)
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"[导出] TXT 已保存：{path}")
    return path


# ---------- Word 导出（python-docx） ----------

def export_to_word(basename: str, title: str, sections: Dict[str, str]) -> str:
    """
    sections: {"任务说明": "...", "Planner 结果": "...", ...}
    """
    try:
        from docx import Document
    except ImportError:
        print("⚠️ 未安装 python-docx，无法导出 Word。请先执行：pip install python-docx")
        return ""

    export_dir = ensure_export_dir()
    filename = f"{basename}_{timestamp_str()}.docx"
    path = os.path.join(export_dir, filename)

    doc = Document()
    doc.add_heading(title, level=1)

    for sec_title, sec_content in sections.items():
        doc.add_heading(sec_title, level=2)
        for line in sec_content.splitlines():
            doc.add_paragraph(line)

    doc.save(path)
    print(f"[导出] Word 已保存：{path}")
    return path


# ---------- PPT 导出（python-pptx） ----------

def export_to_ppt(basename: str, title: str, slides: List[Dict[str, List[str]]]) -> str:
    """
    slides: 每页一个 dict，如：
      {"title": "第1页 标题", "bullets": ["要点1", "要点2", ...]}
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        print("⚠️ 未安装 python-pptx，无法导出 PPT。请先执行：pip install python-pptx")
        return ""

    export_dir = ensure_export_dir()
    filename = f"{basename}_{timestamp_str()}.pptx"
    path = os.path.join(export_dir, filename)

    prs = Presentation()

    # 封面页
    title_slide_layout = prs.slide_layouts[0]  # 标题 + 副标题
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = "由 Graph Agent 自动生成"

    # 内容页
    bullet_layout = prs.slide_layouts[1]  # 标题 + 内容
    for s in slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = s.get("title", "未命名页面")
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        bullets = s.get("bullets", [])
        if not bullets:
            continue
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(20)

    prs.save(path)
    print(f"[导出] PPT 已保存：{path}")
    return path


# ---------- Excel 导出（openpyxl） ----------

def export_to_excel(basename: str, rows: List[List[str]]) -> str:
    """
    rows: 每一行一个 list，例如：
      [["类型", "内容"], ["任务", "xxx"], ["Planner", "..."], ...]
    """
    try:
        from openpyxl import Workbook
    except ImportError:
        print("⚠️ 未安装 openpyxl，无法导出 Excel。请先执行：pip install openpyxl")
        return ""

    export_dir = ensure_export_dir()
    filename = f"{basename}_{timestamp_str()}.xlsx"
    path = os.path.join(export_dir, filename)

    wb = Workbook()
    ws = wb.active
    ws.title = "GraphResult"

    for row in rows:
        ws.append(row)

    wb.save(path)
    print(f"[导出] Excel 已保存：{path}")
    return path
