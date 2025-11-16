# teaching_exporter.py
# -*- coding: utf-8 -*-

from typing import Any, Dict, List
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from openpyxl import Workbook


def export_to_word(result: Dict[str, Any], filename: str) -> None:
    """
    根据 teaching_result 结构导出 Word 教案
    result 结构约定：
    {
        "course_name": str,
        "grade": str,
        "goals": [str, ...],
        "activities": [str, ...],
        "ppt_outline": [{"title": str, "bullets": [str, ...]}, ...],
        "intro": str,
        "resources": [str, ...]
    }
    """
    doc = Document()

    # 标题
    course_name = result.get("course_name", "未命名课程")
    grade = result.get("grade", "")
    doc.add_heading(course_name, level=1)
    if grade:
        doc.add_paragraph(f"适用对象：{grade}")

    # 课程简介
    intro = result.get("intro", "")
    if intro:
        doc.add_heading("一、课程简介", level=2)
        doc.add_paragraph(intro)

    # 教学目标
    goals: List[str] = result.get("goals", [])
    if goals:
        doc.add_heading("二、教学目标", level=2)
        for i, g in enumerate(goals, 1):
            doc.add_paragraph(f"{i}. {g}")

    # 课堂活动
    activities: List[str] = result.get("activities", [])
    if activities:
        doc.add_heading("三、课堂活动设计", level=2)
        for i, act in enumerate(activities, 1):
            p = doc.add_paragraph(f"活动{i}：")
            p.add_run(act)

    # PPT 结构
    ppt_outline: List[Dict[str, Any]] = result.get("ppt_outline", [])
    if ppt_outline:
        doc.add_heading("四、PPT 结构大纲", level=2)
        for i, slide in enumerate(ppt_outline, 1):
            title = slide.get("title", f"第{i}页")
            bullets = slide.get("bullets", [])
            doc.add_paragraph(f"{i}. {title}")
            for b in bullets:
                doc.add_paragraph(f"- {b}", style="List Bullet")

    # 教学资源
    resources: List[str] = result.get("resources", [])
    if resources:
        doc.add_heading("五、教学资源与参考资料", level=2)
        for i, r in enumerate(resources, 1):
            doc.add_paragraph(f"{i}. {r}")

    doc.save(filename)
    print(f"[Export] Word 教案已生成：{filename}")


def export_to_ppt(result: Dict[str, Any], filename: str) -> None:
    """
    根据 teaching_result 结构导出 PPT
    """
    prs = Presentation()

    # 封面页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    course_name = result.get("course_name", "未命名课程")
    grade = result.get("grade", "")
    title.text = course_name
    subtitle.text = f"适用对象：{grade}" if grade else "教学设计与课件示例"

    # 内容大纲页
    ppt_outline: List[Dict[str, Any]] = result.get("ppt_outline", [])
    for slide_info in ppt_outline:
        layout = prs.slide_layouts[1]  # 标题 + 内容
        s = prs.slides.add_slide(layout)
        s_title = s.shapes.title
        body = s.placeholders[1]

        s_title.text = slide_info.get("title", "未命名小节")
        tf = body.text_frame
        tf.clear()

        bullets = slide_info.get("bullets", [])
        for i, b in enumerate(bullets):
            if i == 0:
                tf.text = b
            else:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0

    # 调整字体大小（可选，简单示例）
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(20)

    prs.save(filename)
    print(f"[Export] PPT 课件已生成：{filename}")


def export_to_excel(result: Dict[str, Any], filename: str) -> None:
    """
    根据 teaching_result 结构导出 Excel（主要记录资源和目标）
    """
    wb = Workbook()

    # Sheet1: 课程基本信息 + 教学目标
    ws1 = wb.active
    ws1.title = "课程与目标"

    ws1.append(["字段", "内容"])
    ws1.append(["课程名称", result.get("course_name", "")])
    ws1.append(["适用对象", result.get("grade", "")])
    ws1.append([])

    ws1.append(["教学目标编号", "教学目标内容"])
    goals = result.get("goals", [])
    for i, g in enumerate(goals, 1):
        ws1.append([i, g])

    # Sheet2: 教学资源
    ws2 = wb.create_sheet(title="教学资源")
    ws2.append(["序号", "资源说明"])
    resources = result.get("resources", [])
    for i, r in enumerate(resources, 1):
        ws2.append([i, r])

    wb.save(filename)
    print(f"[Export] Excel 资源表已生成：{filename}")


# 可选：独立测试本模块
def build_sample_teaching_result() -> Dict[str, Any]:
    return {
        "course_name": "人工智能基础入门",
        "grade": "中职一年级 / 大专一年级",
        "goals": [
            "理解人工智能的基本概念与发展脉络",
            "能说出至少三类典型的 AI 应用场景",
            "初步形成用 AI 辅助学习与解决问题的意识",
        ],
        "activities": [
            "导入讨论：你身边遇到过哪些“看起来像 AI” 的应用？",
            "小组任务：用自己的话解释什么是人工智能，并举 1 个例子。",
            "应用头脑风暴：本专业将来哪些岗位可能会用到 AI？",
        ],
        "ppt_outline": [
            {
                "title": "课程导入与学习目标",
                "bullets": ["课程背景与意义", "本节课学习目标", "课堂活动预告"],
            },
            {
                "title": "什么是人工智能",
                "bullets": ["直观例子：语音助手、推荐系统等", "简单概念解释", "常见误解澄清"],
            },
        ],
        "intro": "本课程面向零基础学生，通过实例与案例引导学生理解人工智能的基本概念、典型应用与发展趋势，为后续深入学习打下基础。",
        "resources": [
            "视频：李宏毅《机器学习》导论第一讲（截取片段）",
            "文章：通俗解读——什么是人工智能？",
            "工具：可视化机器学习体验网站（如 Teachable Machine 等）",
        ],
    }


if __name__ == "__main__":
    # 简单自测
    sample = build_sample_teaching_result()
    export_to_word(sample, "demo_教案.docx")
    export_to_ppt(sample, "demo_课件.pptx")
    export_to_excel(sample, "demo_资源表.xlsx")
    print("[Export] demo 导出完成。")
