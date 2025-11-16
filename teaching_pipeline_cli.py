# teaching_pipeline_cli.py
"""
教学流水线核心逻辑（CLI + Web 共用）

功能：
1. 调用 OpenAI 生成教学设计（目标 / 大纲 / 活动 / 作业 / 反思）
2. 简单 Memory（GraphMemory），记录最近若干次任务
3. 导出 Word / PPT / Excel

供 teaching_pipeline_web.py 导入使用：
- run_teaching_pipeline
- export_to_word
- export_to_ppt
- export_to_excel
- GraphMemory
"""

import os
import json
from dataclasses import dataclass, field
from datetime import datetime
from typing import List, Dict, Any, Optional

try:
    import streamlit as st
except ImportError:
    st = None  # 本地 CLI 模式可能没有 streamlit

from openai import OpenAI, AuthenticationError, OpenAIError

# 默认模型名称 —— 注意：不要写不存在的模型
DEFAULT_MODEL = "gpt-4.1-mini"

# =========================
#  一些简单的工具函数
# =========================


def _get_openai_api_key() -> str:
    """
    从环境变量或 Streamlit secrets 中获取 OPENAI_API_KEY。
    """
    api_key = os.getenv("OPENAI_API_KEY")

    if (not api_key) and st is not None:
        try:
            api_key = st.secrets.get("OPENAI_API_KEY", None)
        except Exception:
            api_key = None

    if not api_key:
        # 这里抛 RuntimeError，让上层统一处理成友好提示
        raise RuntimeError(
            "未找到 OPENAI_API_KEY。\n"
            "请在本地系统环境变量中设置，或在 Streamlit Cloud 的 Secrets 中配置：\n\n"
            "[default]\nOPENAI_API_KEY = \"sk-xxxx\"\n"
        )
    return api_key


def _get_openai_client() -> OpenAI:
    api_key = _get_openai_api_key()
    return OpenAI(api_key=api_key)


def call_llm(system_prompt: str, user_content: str, model: str = DEFAULT_MODEL) -> str:
    """
    统一的 LLM 调用封装。
    返回纯文本字符串。
    """
    client = _get_openai_client()

    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content},
            ],
            temperature=0.4,
        )
        return (resp.choices[0].message.content or "").strip()

    except AuthenticationError as e:
        raise RuntimeError(
            "OpenAI 身份验证失败：请检查 OPENAI_API_KEY 是否正确，"
            "并确认该账户具有 API 使用权限。"
        ) from e

    except OpenAIError as e:
        raise RuntimeError(f"OpenAI 调用失败：{e.__class__.__name__}：{e}") from e


# =========================
#  Memory：GraphMemory
# =========================


@dataclass
class MemoryRecord:
    timestamp: str
    task: str
    course_title: str
    student_level: str
    duration: int
    summary: str


@dataclass
class GraphMemory:
    """
    非严格意义的“图记忆”，这里做成简单队列形式即可。
    """

    max_len: int = 10
    records: List[MemoryRecord] = field(default_factory=list)

    # 兼容多种写入方法：add / append / write
    def add_record(self, record: MemoryRecord) -> None:
        self.records.append(record)
        if len(self.records) > self.max_len:
            self.records = self.records[-self.max_len :]

    def add(self, record: MemoryRecord) -> None:
        self.add_record(record)

    def append(self, record: MemoryRecord) -> None:
        self.add_record(record)

    def write(self, record: MemoryRecord) -> None:
        self.add_record(record)

    def to_json(self) -> str:
        return json.dumps([r.__dict__ for r in self.records], ensure_ascii=False, indent=2)

    @classmethod
    def from_json(cls, s: str, max_len: int = 10) -> "GraphMemory":
        try:
            data = json.loads(s)
        except Exception:
            return cls(max_len=max_len)

        records = []
        for item in data:
            try:
                records.append(
                    MemoryRecord(
                        timestamp=item.get("timestamp", ""),
                        task=item.get("task", ""),
                        course_title=item.get("course_title", ""),
                        student_level=item.get("student_level", ""),
                        duration=item.get("duration", 0),
                        summary=item.get("summary", ""),
                    )
                )
            except Exception:
                continue
        mem = cls(max_len=max_len, records=records)
        return mem


# =========================
#  教学流水线主逻辑
# =========================


def _build_system_prompt() -> str:
    return (
        "你是一名经验丰富的教学设计专家，擅长为中小学和高校教师设计结构化、可落地的教学方案。"
        "请根据用户给出的课程标题、学习者水平、课时长度和教学任务，生成完整的教学设计，"
        "包括：三维目标 / 知识结构 / 教学环节 / 课堂活动 / 作业设计 / 教学反思建议等。"
    )


def _build_user_prompt(task: str, course_title: str, student_level: str, duration: int) -> str:
    return f"""
【课程标题】
{course_title}

【学生水平】
{student_level}

【单节课时长】
约 {duration} 分钟

【本次教学任务】
{task}

请输出结构化的教学设计，建议使用「分条」形式，包含但不限于：
1. 教学目标（知识与能力 / 过程与方法 / 情感态度与价值观）
2. 学情分析（简要）
3. 教学重点与难点
4. 教学思路与整体结构
5. 教学过程设计（导入 / 新授 / 活动 / 巩固 / 小结）
6. 课堂活动与问题设计（尽量具体，可以写成步骤）
7. 课后作业与延伸任务
8. 教学反思建议（给教师的改进提示）

要求：
- 语言简洁明了，但要足够具体，便于教师直接使用或稍作修改后使用；
- 尽量加入互动、讨论、分组等活动元素；
- 不要输出与技术实现或代码有关的内容。
""".strip()


def _parse_plan_to_structured_dict(raw_text: str) -> Dict[str, Any]:
    """
    简单解析 LLM 返回的教学设计文本，拆分为若干字段。
    为了稳健，不做复杂正则，只做大概切分。
    """
    result: Dict[str, Any] = {
        "goals": "",
        "analysis": "",
        "key_points": "",
        "process": "",
        "activities": "",
        "homework": "",
        "reflection": "",
        "raw_plan_text": raw_text,
    }

    # 按行拆分
    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    current_key = None
    buffer: Dict[str, List[str]] = {
        "goals": [],
        "analysis": [],
        "key_points": [],
        "process": [],
        "activities": [],
        "homework": [],
        "reflection": [],
    }

    def _set_key_from_title(line: str) -> Optional[str]:
        mapping = {
            "目标": "goals",
            "教学目标": "goals",
            "学情分析": "analysis",
            "学生分析": "analysis",
            "教学重点": "key_points",
            "重点与难点": "key_points",
            "教学难点": "key_points",
            "教学过程": "process",
            "教学环节": "process",
            "课堂活动": "activities",
            "活动设计": "activities",
            "作业": "homework",
            "课后作业": "homework",
            "教学反思": "reflection",
            "反思建议": "reflection",
        }
        for k, v in mapping.items():
            if k in line:
                return v
        return None

    for line in lines:
        new_key = _set_key_from_title(line)
        if new_key is not None:
            current_key = new_key
            continue

        if current_key is None:
            # 还没识别到章节标题，先丢到 goals 里
            buffer["goals"].append(line)
        else:
            buffer[current_key].append(line)

    for k in buffer:
        result[k] = "\n".join(buffer[k]).strip()

    return result


def _build_ppt_outline(plan_struct: Dict[str, Any]) -> str:
    """
    根据教学设计生成 PPT 大纲文本（纯文本），供后续导出 PPT 使用。
    """
    goals = plan_struct.get("goals", "")
    process = plan_struct.get("process", "")
    activities = plan_struct.get("activities", "")
    homework = plan_struct.get("homework", "")

    parts = [
        "封面：课程标题 / 授课人 / 时间",
        "",
        "目录：",
        "1. 教学目标",
        "2. 学习准备与导入",
        "3. 知识讲解与示例",
        "4. 课堂活动与练习",
        "5. 小结与作业",
        "",
        "【详细内容建议】",
        "",
        "一、教学目标（可直接放在第 1 页正文）",
        goals or "（此处填写教学目标）",
        "",
        "二、教学过程要点（用于生成 3-4 页讲解型幻灯片）",
        process or "（此处填写教学过程梗概）",
        "",
        "三、课堂活动（可单独 1-2 页展示活动步骤）",
        activities or "（此处填写课堂活动步骤）",
        "",
        "四、课后作业（最后一页）",
        homework or "（此处填写作业或延伸任务）",
    ]

    return "\n".join(parts).strip()


def run_teaching_pipeline(
    task: str,
    course_title: str,
    student_level: str,
    duration: int,
    memory: Optional[GraphMemory] = None,
    model: str = DEFAULT_MODEL,
) -> Dict[str, Any]:
    """
    教学流水线主入口。

    返回 dict，包含：
    - task / course_title / student_level / duration
    - goals / analysis / key_points / process / activities / homework / reflection
    - ppt_outline
    - raw_plan_text / raw_ppt_text
    """
    system_prompt = _build_system_prompt()
    user_prompt = _build_user_prompt(task, course_title, student_level, duration)

    raw_plan = call_llm(system_prompt, user_prompt, model=model)
    plan_struct = _parse_plan_to_structured_dict(raw_plan)
    ppt_outline = _build_ppt_outline(plan_struct)

    result: Dict[str, Any] = {
        "task": task,
        "course_title": course_title,
        "student_level": student_level,
        "duration": duration,
        "model": model,
        # 结构化字段
        "goals": plan_struct.get("goals", ""),
        "analysis": plan_struct.get("analysis", ""),
        "key_points": plan_struct.get("key_points", ""),
        "process": plan_struct.get("process", ""),
        "activities": plan_struct.get("activities", ""),
        "homework": plan_struct.get("homework", ""),
        "reflection": plan_struct.get("reflection", ""),
        # 文本原文
        "raw_plan_text": plan_struct.get("raw_plan_text", raw_plan),
        "ppt_outline": ppt_outline,
        "raw_ppt_text": ppt_outline,
    }

    # 写入 Memory（如果传入了）
    if memory is not None:
        try:
            summary = f"{course_title} / {student_level} / {duration} 分钟 / 任务：{task[:50]}..."
            record = MemoryRecord(
                timestamp=datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                task=task,
                course_title=course_title,
                student_level=student_level,
                duration=duration,
                summary=summary,
            )
            # 兼容多种方法名
            if hasattr(memory, "add_record"):
                memory.add_record(record)
            elif hasattr(memory, "add"):
                memory.add(record)
            elif hasattr(memory, "append"):
                memory.append(record)
            elif hasattr(memory, "write"):
                memory.write(record)
            else:
                print("GraphMemory 中未发现可用的写入方法，已跳过记忆写入。")
        except Exception as e:
            print(f"写入 GraphMemory 时出现异常：{e}")

    return result


# =========================
#  导出 Word / PPT / Excel
# =========================


def export_to_word(result: Dict[str, Any], filepath: str) -> None:
    """
    将教学设计导出为 Word（.docx）。
    """
    from docx import Document

    doc = Document()

    title = result.get("course_title", "教学设计")
    doc.add_heading(title, level=1)

    meta_p = doc.add_paragraph()
    meta_p.add_run("学生水平：").bold = True
    meta_p.add_run(result.get("student_level", ""))
    meta_p.add_run("    课时：").bold = True
    meta_p.add_run(str(result.get("duration", "")) + " 分钟")

    doc.add_heading("一、教学目标", level=2)
    doc.add_paragraph(result.get("goals", "") or "（此处填写教学目标）")

    doc.add_heading("二、学情分析", level=2)
    doc.add_paragraph(result.get("analysis", "") or "（此处填写学情分析）")

    doc.add_heading("三、教学重点与难点", level=2)
    doc.add_paragraph(result.get("key_points", "") or "（此处填写教学重点与难点）")

    doc.add_heading("四、教学过程", level=2)
    doc.add_paragraph(result.get("process", "") or "（此处填写教学过程）")

    doc.add_heading("五、课堂活动设计", level=2)
    doc.add_paragraph(result.get("activities", "") or "（此处填写课堂活动）")

    doc.add_heading("六、课后作业与延伸", level=2)
    doc.add_paragraph(result.get("homework", "") or "（此处填写作业与延伸）")

    doc.add_heading("七、教学反思建议", level=2)
    doc.add_paragraph(result.get("reflection", "") or "（此处填写教学反思）")

    doc.save(filepath)


def export_to_ppt(result: Dict[str, Any], filepath: str) -> None:
    """
    将教学设计导出为 PPTX。
    注意：TextFrame.add_paragraph() 不带参数，先 add_paragraph() 再设置 p.text。
    """
    from pptx import Presentation

    prs = Presentation()

    course_title = result.get("course_title", "教学课件")
    student_level = result.get("student_level", "")
    duration = result.get("duration", "")

    # 封面
    slide_layout = prs.slide_layouts[0]  # 标题+副标题
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = course_title
    subtitle = slide.placeholders[1]
    subtitle.text = f"{student_level} / 约 {duration} 分钟"

    # 教学目标
    slide_layout = prs.slide_layouts[1]  # 标题+正文
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "教学目标"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    goals = result.get("goals", "") or "（此处填写教学目标）"
    for line in goals.splitlines():
        if not line.strip():
            continue
        if not body.text:
            body.text = line.strip()
        else:
            p = body.add_paragraph()
            p.text = line.strip()

    # 教学过程 / 知识讲解
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "教学过程要点"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    process = result.get("process", "") or "（此处填写教学过程要点）"
    for line in process.splitlines():
        if not line.strip():
            continue
        if not body.text:
            body.text = line.strip()
        else:
            p = body.add_paragraph()
            p.text = line.strip()

    # 课堂活动
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "课堂活动"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    activities = result.get("activities", "") or "（此处填写课堂活动步骤）"
    for line in activities.splitlines():
        if not line.strip():
            continue
        if not body.text:
            body.text = line.strip()
        else:
            p = body.add_paragraph()
            p.text = line.strip()

    # 作业
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "课后作业与延伸"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    homework = result.get("homework", "") or "（此处填写作业与延伸任务）"
    for line in homework.splitlines():
        if not line.strip():
            continue
        if not body.text:
            body.text = line.strip()
        else:
            p = body.add_paragraph()
            p.text = line.strip()

    prs.save(filepath)


def export_to_excel(result: Dict[str, Any], filepath: str) -> None:
    """
    将教学设计导出为 Excel（.xlsx）。
    """
    from openpyxl import Workbook

    wb = Workbook()

    # Sheet1：整体教学设计
    ws1 = wb.active
    ws1.title = "教学设计"

    ws1["A1"] = "课程标题"
    ws1["B1"] = result.get("course_title", "")
    ws1["A2"] = "学生水平"
    ws1["B2"] = result.get("student_level", "")
    ws1["A3"] = "课时（分钟）"
    ws1["B3"] = result.get("duration", "")

    ws1["A5"] = "教学目标"
    ws1["B5"] = result.get("goals", "")

    ws1["A6"] = "学情分析"
    ws1["B6"] = result.get("analysis", "")

    ws1["A7"] = "重点与难点"
    ws1["B7"] = result.get("key_points", "")

    ws1["A8"] = "教学过程"
    ws1["B8"] = result.get("process", "")

    ws1["A9"] = "课堂活动"
    ws1["B9"] = result.get("activities", "")

    ws1["A10"] = "课后作业"
    ws1["B10"] = result.get("homework", "")

    ws1["A11"] = "教学反思"
    ws1["B11"] = result.get("reflection", "")

    # Sheet2：PPT 大纲
    ws2 = wb.create_sheet(title="PPT大纲")
    ws2["A1"] = "PPT 大纲"
    ws2["A2"] = result.get("ppt_outline", "")

    # Sheet3：原始文本
    ws3 = wb.create_sheet(title="原始文本")
    ws3["A1"] = "教学设计原文"
    ws3["A2"] = result.get("raw_plan_text", "")

    wb.save(filepath)


# =========================
#  可选：本地 CLI 测试入口
# =========================

if __name__ == "__main__":
    # 简单命令行测试：python teaching_pipeline_cli.py
    print("=== 教学流水线 CLI 测试 ===")
    task = input("请输入本次教学任务：").strip() or "讲解一次函数的图像与性质"
    course_title = input("课程标题：").strip() or "一次函数基础"
    student_level = input("学生水平：").strip() or "初二学生"
    duration_str = input("课时（分钟）：").strip() or "40"
    duration = int(duration_str)

    memory = GraphMemory(max_len=10)
    result = run_teaching_pipeline(
        task=task,
        course_title=course_title,
        student_level=student_level,
        duration=duration,
        memory=memory,
    )

    print("\n=== 本次 AI 教学设计结果（摘要） ===")
    print("教学目标：")
    print(result.get("goals", "")[:500])
    print("\n教学过程：")
    print(result.get("process", "")[:500])

    # 导出示例
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    word_path = f"teaching_pipeline_{ts}.docx"
    ppt_path = f"teaching_pipeline_{ts}.pptx"
    excel_path = f"teaching_pipeline_{ts}.xlsx"

    export_to_word(result, word_path)
    export_to_ppt(result, ppt_path)
    export_to_excel(result, excel_path)

    print(f"\n已导出 Word / PPT / Excel：\n{word_path}\n{ppt_path}\n{excel_path}")
