# -*- coding: utf-8 -*-
"""
教学设计结果 → 导出为 Word / PPT / Excel 的简单示例

文件名建议：teaching_export_demo.py
运行方式：python teaching_export_demo.py
"""

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from openpyxl import Workbook


# ======================================================
# 1. 定义一个统一的数据结构（教学结果 teaching_result）
# ======================================================

def build_sample_teaching_result():
    """
    这里先手工构造一个示例教学结果。
    以后可以替换成：从 Graph / Agent 的最终输出整理而来。
    """
    teaching_result = {
        "course_name": "人工智能基础入门",
        "grade": "中职一年级",

        # 3 条教学目标
        "goals": [
            "理解人工智能的基本概念和主要发展阶段，能够用自己的话进行简要说明。",
            "能举例说明人工智能在日常生活和行业中的典型应用场景，增强学习兴趣。",
            "初步具备使用简单 AI 工具完成学习任务的意识与意愿，形成积极的学习态度。",
        ],

        # 课堂活动（这里直接用一大段文字，你可以替换为真实 Agent 输出）
        "activities": """一、活动一：课堂导入讨论
【目标】激活学生已有经验，了解他们对“人工智能”的直观认识。
【步骤】
1. 教师展示几张与 AI 应用相关的图片或短视频；
2. 学生分组讨论“哪些是人工智能？为什么？”（3–5 分钟）；
3. 小组代表分享观点，教师进行归纳整理。
【材料】多媒体设备、图片或视频、白板或黑板。

二、活动二：用自己的话解释“什么是人工智能”
【目标】让学生尝试用通俗语言概括人工智能的核心概念，培养表达能力。
【步骤】
1. 每组领取一张任务纸，上面写着“用 1–2 句话解释什么是人工智能”；
2. 小组讨论并写出自己的定义或类比（10 分钟）；
3. 全班展示，小组互评，教师给出专业版概括。
【材料】任务纸、签字笔、投影或黑板。

三、活动三：AI 应用场景头脑风暴
【目标】帮助学生意识到 AI 与自己生活和未来职业之间的关系。
【步骤】
1. 教师示范 2–3 个典型场景（如智能客服、人脸识别等）；
2. 学生分组头脑风暴，写出 3–5 个他们能想到的 AI 应用场景（10 分钟）；
3. 小组选择一个最有兴趣的场景进行简要说明：是什么？解决了什么问题？
【材料】便签纸、大张白纸或白板、彩笔。""",

        # PPT 结构（8 页）
        "ppt_outline": [
            {
                "title": "课程导入与学习目标",
                "bullets": [
                    "课程背景：为什么要学人工智能？",
                    "本课程面向的学生与适用场景",
                    "本节课的学习目标（知识 / 能力 / 态度）"
                ]
            },
            {
                "title": "人工智能基本概念",
                "bullets": [
                    "什么是“智能”？什么是“人工智能”？",
                    "图灵测试与机器模拟人类智能的早期设想",
                    "弱人工智能与强人工智能的简单区分"
                ]
            },
            {
                "title": "人工智能发展简史",
                "bullets": [
                    "从“符号主义 AI”到“机器学习”的转变",
                    "深度学习兴起与算力、数据的推动",
                    "当代大模型（如 ChatGPT）的出现"
                ]
            },
            {
                "title": "人工智能典型应用场景",
                "bullets": [
                    "生活中的 AI：推荐系统、语音助手、导航等",
                    "行业中的 AI：制造、金融、医疗、教育等",
                    "简单思考：哪些地方还可以用 AI？"
                ]
            },
            {
                "title": "课堂活动一：导入讨论",
                "bullets": [
                    "展示图片 / 视频，引导学生说出直观感受",
                    "小组讨论：什么是人工智能？",
                    "教师总结学生的观点"
                ]
            },
            {
                "title": "课堂活动二：我的 AI 定义",
                "bullets": [
                    "小组任务：用 1–2 句话写出你心中的 AI 定义",
                    "小组展示与互评",
                    "教师给出更专业、更严谨的版本"
                ]
            },
            {
                "title": "课堂活动三：AI 应用头脑风暴",
                "bullets": [
                    "头脑风暴规则与分组方式说明",
                    "学生分组写出生活 / 职业中的 AI 场景",
                    "小组分享与教师点评"
                ]
            },
            {
                "title": "课堂小结与课后延伸",
                "bullets": [
                    "本节课的三个关键词回顾",
                    "学生今天获得的一个新认识 / 一个新问题",
                    "课后可看的视频 / 可体验的 AI 工具推荐"
                ]
            }
        ],

        # 课程简介
        "intro": (
            "本课程面向中职一年级学生，围绕“什么是人工智能、人工智能从哪里来、"
            "它正在改变什么”这三个核心问题展开。通过通俗讲解、课堂讨论和应用场景"
            "分析，帮助学生建立对人工智能的初步整体认识，激发学习兴趣，为后续进一"
            "步学习 AI 技术和工具打下基础。"
        ),

        # 教学资源
        "resources": [
            {
                "type": "教材",
                "name": "《人工智能：一种现代的方法》（第 3 版）",
                "note": "Stuart Russell & Peter Norvig，作为教师参考，不要求学生细读。"
            },
            {
                "type": "参考书",
                "name": "《机器学习》（周志华）",
                "note": "可作为后续进阶阅读材料。"
            },
            {
                "type": "在线课程",
                "name": "Coursera – Machine Learning（吴恩达）",
                "note": "配套视频，可选修。"
            },
            {
                "type": "工具网站",
                "name": "ChatGPT / 通义千问 / 讯飞星火等在线对话式 AI",
                "note": "用于课堂演示和学生课后体验。"
            }
        ]
    }
    return teaching_result


# ======================================================
# 2. 导出 Word 教案
# ======================================================

def export_to_word(result: dict, filename: str = "教案_人工智能基础入门.docx"):
    doc = Document()

    # 设置中文默认字体（不是必须，但有助于显示效果）
    style = doc.styles['Normal']
    style.font.name = '宋体'
    style._element.rPr.rFonts.set(qn('eastAsia'), '宋体') if 'eastAsia' in dir(style._element.rPr.rFonts) else None

    # 标题
    doc.add_heading(f"{result['course_name']} 教学设计方案", level=1)

    # 基本信息
    p = doc.add_paragraph()
    run = p.add_run(f"适用对象：{result.get('grade', '未标明')}")
    run.bold = True

    # 一、课程简介
    doc.add_heading("一、课程简介", level=2)
    doc.add_paragraph(result.get("intro", ""))

    # 二、教学目标
    doc.add_heading("二、教学目标", level=2)
    goals = result.get("goals", [])
    for i, g in enumerate(goals, start=1):
        doc.add_paragraph(f"{i}. {g}", style="List Number")

    # 三、课堂活动设计
    doc.add_heading("三、课堂活动设计", level=2)
    doc.add_paragraph(result.get("activities", ""))

    # 四、PPT 结构大纲
    doc.add_heading("四、PPT 结构大纲", level=2)
    ppt_outline = result.get("ppt_outline", [])
    for i, slide in enumerate(ppt_outline, start=1):
        doc.add_paragraph(f"第 {i} 页：{slide['title']}", style="List Bullet")
        for b in slide.get("bullets", []):
            doc.add_paragraph(f"  - {b}", style="List Bullet 2")

    # 五、教学资源
    doc.add_heading("五、教学资源", level=2)
    resources = result.get("resources", [])
    for r in resources:
        line = f"{r.get('type', '')}：{r.get('name', '')}"
        if r.get("note"):
            line += f"（备注：{r['note']}）"
        doc.add_paragraph(line, style="List Bullet")

    doc.save(filename)
    print(f"[OK] 已生成 Word 教案：{filename}")


# ======================================================
# 3. 导出 PPT 结构
# ======================================================

def export_to_ppt(result: dict, filename: str = "课件_人工智能基础入门.pptx"):
    prs = Presentation()

    course_name = result["course_name"]
    ppt_outline = result.get("ppt_outline", [])

    # 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # 标题 + 副标题
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = course_name
    subtitle.text = "教学设计与课堂活动大纲"

    # 内容页
    for slide_info in ppt_outline:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题 + 内容
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = slide_info["title"]

        tf = body.text_frame
        tf.clear()  # 清空默认文本

        first = True
        for b in slide_info.get("bullets", []):
            if first:
                tf.text = b
                first = False
            else:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0

    prs.save(filename)
    print(f"[OK] 已生成 PPT 大纲：{filename}")


# ======================================================
# 4. 导出 Excel 资源清单
# ======================================================

def export_to_excel(result: dict, filename: str = "资源清单_人工智能基础入门.xlsx"):
    wb = Workbook()
    ws = wb.active
    ws.title = "资源清单"

    # 表头
    ws.append(["类型", "名称", "备注"])

    # 内容
    for r in result.get("resources", []):
        ws.append([
            r.get("type", ""),
            r.get("name", ""),
            r.get("note", "")
        ])

    wb.save(filename)
    print(f"[OK] 已生成 Excel 资源表：{filename}")


# ======================================================
# 5. 主函数：统一测试三个导出功能
# ======================================================

def main():
    # 1）构造示例教学结果
    teaching_result = build_sample_teaching_result()

    # 2）导出为 Word / PPT / Excel
    export_to_word(teaching_result, "教案_人工智能基础入门.docx")
    export_to_ppt(teaching_result, "课件_人工智能基础入门.pptx")
    export_to_excel(teaching_result, "资源清单_人工智能基础入门.xlsx")

    print("\n=== 全部导出完成，请到当前目录查看生成的 3 个文件 ===")


if __name__ == "__main__":
    main()
